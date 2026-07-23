"""File-based document upload and text extraction.

Accepts PDF, DOCX, and TXT files via multipart upload.
Extracts text content and stores in document_uploads table.
"""
from __future__ import annotations

import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)


import re


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF bytes using pdfplumber (preferred) or pypdf."""
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            pages = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
    except ImportError:
        pass

    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        pass

    raise ImportError("No PDF library available. Install pdfplumber or pypdf.")


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except ImportError:
        # Fallback: docx files are ZIP archives with XML
        import zipfile
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            with zf.open("word/document.xml") as f:
                tree = ET.parse(f)

        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        paragraphs = []
        for p in tree.iter(f"{{{ns['w']}}}p"):
            texts = []
            for t in p.iter(f"{{{ns['w']}}}t"):
                if t.text:
                    texts.append(t.text)
            if texts:
                paragraphs.append("".join(texts))
        return "\n\n".join(paragraphs)


def extract_text_from_html(file_bytes: bytes) -> str:
    """Gate A (2026-07-16): .html/.htm through the same main-text
    extractor the web lane ships (chrome-stripped, title recovered as
    the first line so detection and chunk labels see it)."""
    from ..search.fetch import extract_main_text

    title, body = extract_main_text(
        file_bytes.decode("utf-8", errors="replace")
    )
    return (f"{title.strip()}\n\n{body}" if (title or "").strip()
            else body)


def extract_transcript_from_subtitles(file_bytes: bytes) -> str:
    """Gate A (2026-07-16): .vtt/.srt -> speaker-attributed transcript
    text. Zoom/Meet exports carry 'Name: text' in cues (or <v Name>
    voice tags); WEBVTT headers, NOTE/STYLE blocks, cue ids, and
    timestamp lines are dropped. The result lands on the transcript
    detected_type — the dynamics profile for free."""
    text = file_bytes.decode("utf-8", errors="replace")
    lines: list[str] = []
    for rawline in text.splitlines():
        s = rawline.strip().lstrip("\ufeff")
        if not s:
            continue
        upper = s.upper()
        if upper.startswith(("WEBVTT", "NOTE", "STYLE", "REGION")):
            continue
        if "-->" in s:
            continue
        if s.isdigit():
            continue
        m = re.match(r"<v\s+([^>]+)>(.*)", s)
        if m:
            s = f"{m.group(1).strip()}: {m.group(2)}"
        s = re.sub(r"<[^>]+>", "", s).strip()
        if s:
            lines.append(s)
    return "\n".join(lines)


def extract_text_from_file(
    file_bytes: bytes,
    filename: str,
    mime: Optional[str] = None,
) -> str:
    """Extract text from a file: extension dispatch first, declared
    MIME as the fallback for extensionless sources (C3, wired by
    Gate H)."""
    lower = filename.lower()

    if lower.endswith(".eml"):
        return extract_chat_from_eml(file_bytes)
    elif lower.endswith(".mbox"):
        return extract_chat_from_mbox(file_bytes)
    elif lower.endswith(".xlsx"):
        return extract_tabular_from_xlsx(file_bytes)
    elif lower.endswith(".csv"):
        return extract_tabular_from_delimited(file_bytes, ",")
    elif lower.endswith(".tsv"):
        return extract_tabular_from_delimited(file_bytes, "\t")
    elif lower.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif lower.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif lower.endswith(".html") or lower.endswith(".htm"):
        return extract_text_from_html(file_bytes)
    elif lower.endswith(".vtt") or lower.endswith(".srt"):
        return extract_transcript_from_subtitles(file_bytes)
    elif lower.endswith(".txt") or lower.endswith(".md"):
        return file_bytes.decode("utf-8", errors="replace")
    elif lower.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif lower.endswith(".rtf"):
        return extract_text_from_rtf(file_bytes)
    elif lower.endswith(".odt"):
        return extract_text_from_odt(file_bytes)
    elif lower.endswith(".epub"):
        return extract_text_from_epub(file_bytes)
    elif lower.endswith(".ipynb"):
        return extract_text_from_ipynb(file_bytes)
    else:
        # C3 MIME fallback (wired by Gate H): no recognized extension —
        # map the declared MIME to an extension and re-dispatch ONCE
        # (the mapped name always has a known extension, so this cannot
        # recurse further). Serves the connector envelope's no-filename
        # case.
        ext = _MIME_EXTENSIONS.get(
            (mime or "").split(";")[0].strip().lower()
        )
        if ext:
            return extract_text_from_file(file_bytes, f"file{ext}")
        # Try as plain text
        try:
            return file_bytes.decode("utf-8", errors="replace")
        except Exception:
            raise ValueError(f"Unsupported file type: {filename}")


# --- Gate H (2026-07-23): text-adapter batch --------------------------------
# H-Q1=A: no fragment carve for any of these — prose-class documents;
# pptx slides ride as in-text `=== SLIDE N ===` locator markers (decks
# have no set-in-stone format; the general chunker builds the shape).
# H-Q2=A: python-pptx + striprtf in core deps, stdlib fallbacks coded
# (the docx precedent). odt/epub/ipynb are stdlib by design.

PPTX_SLIDE_MARKER = "=== SLIDE "

_MIME_EXTENSIONS = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.oasis.opendocument.text": ".odt",
    "application/rtf": ".rtf",
    "text/rtf": ".rtf",
    "application/epub+zip": ".epub",
    "application/x-ipynb+json": ".ipynb",
    "application/json": ".json",
    "text/html": ".html",
    "text/csv": ".csv",
    "text/tab-separated-values": ".tsv",
    "text/markdown": ".md",
    "message/rfc822": ".eml",
}


def _extract_pptx_stdlib(file_bytes: bytes) -> str:
    """Fallback: pptx is a zip; slide text lives in <a:t> runs inside
    ppt/slides/slideN.xml. Loses tables-as-structure and notes; keeps
    every visible text run in slide order."""
    import zipfile
    import xml.etree.ElementTree as ET

    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    parts: list[str] = []
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        slide_names = sorted(
            (n for n in zf.namelist()
             if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
            key=lambda n: int(re.search(r"(\d+)", n).group(1)),
        )
        for name in slide_names:
            num = int(re.search(r"(\d+)", name).group(1))
            tree = ET.parse(zf.open(name))
            runs = [t.text for t in tree.iter(f"{a_ns}t") if t.text]
            body = "\n".join(r for r in runs if r.strip())
            parts.append(f"{PPTX_SLIDE_MARKER}{num} ===\n{body}".rstrip())
    return "\n\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Slides in order with `=== SLIDE N ===` markers. python-pptx
    reads shapes, tables, and speaker notes; the stdlib fallback keeps
    text runs only."""
    try:
        from pptx import Presentation
    except ImportError:
        return _extract_pptx_stdlib(file_bytes)
    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append("\t".join(cells))
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        body = "\n".join(lines)
        if notes:
            body = f"{body}\nNotes: {notes}" if body else f"Notes: {notes}"
        parts.append(f"{PPTX_SLIDE_MARKER}{i} ===\n{body}".rstrip())
    return "\n\n".join(parts)


_RTF_CONTROL = re.compile(
    r"\\\'[0-9a-fA-F]{2}|\\[a-zA-Z]+-?\d*[ ]?|[{}]|\\[^a-zA-Z]"
)


def _extract_rtf_stdlib(file_bytes: bytes) -> str:
    """Fallback: strip control words/groups. Approximate by design —
    plain documents come through clean; embedded objects degrade."""
    text = file_bytes.decode("latin-1", errors="replace")
    for group in ("fonttbl", "colortbl", "stylesheet", "pict", "info"):
        text = re.sub(
            r"\{\\" + group + r".*?\}", "", text, flags=re.DOTALL,
        )
    text = text.replace("\\par", "\n").replace("\\line", "\n")
    text = _RTF_CONTROL.sub("", text)
    lines = [ln.strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln)


def extract_text_from_rtf(file_bytes: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return _extract_rtf_stdlib(file_bytes)
    return rtf_to_text(
        file_bytes.decode("latin-1", errors="replace"),
    ).strip()


def extract_text_from_odt(file_bytes: bytes) -> str:
    """Stdlib only by design: odt is a zip whose content.xml carries
    every paragraph/heading in <text:p>/<text:h> — a library adds
    nothing for text extraction."""
    import zipfile
    import xml.etree.ElementTree as ET

    t_ns = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        tree = ET.parse(zf.open("content.xml"))
    paragraphs: list[str] = []
    for el in tree.iter():
        if el.tag in (f"{t_ns}p", f"{t_ns}h"):
            text = "".join(el.itertext()).strip()
            if text:
                paragraphs.append(text)
    return "\n\n".join(paragraphs)


def extract_text_from_epub(file_bytes: bytes) -> str:
    """Stdlib zip walk: container.xml -> OPF -> spine order -> each
    xhtml chapter through the SAME html extractor the web lane ships
    (chrome-stripped, title recovered). Chapters join in reading
    order."""
    import zipfile
    import xml.etree.ElementTree as ET

    c_ns = "{urn:oasis:names:tc:opendocument:xmlns:container}"
    o_ns = "{http://www.idpf.org/2007/opf}"
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        container = ET.parse(zf.open("META-INF/container.xml"))
        rootfile = container.find(f".//{c_ns}rootfile")
        opf_path = rootfile.get("full-path")
        opf = ET.parse(zf.open(opf_path))
        base = opf_path.rsplit("/", 1)[0] + "/" if "/" in opf_path else ""
        items = {
            it.get("id"): it.get("href")
            for it in opf.iter(f"{o_ns}item")
        }
        chapters: list[str] = []
        for ref in opf.iter(f"{o_ns}itemref"):
            href = items.get(ref.get("idref"))
            if not href or not href.lower().endswith(
                (".xhtml", ".html", ".htm")
            ):
                continue
            try:
                raw = zf.read(f"{base}{href}")
            except KeyError:
                continue
            text = extract_text_from_html(raw).strip()
            if text:
                chapters.append(text)
    return "\n\n".join(chapters)


def extract_text_from_ipynb(file_bytes: bytes) -> str:
    """Stdlib json: markdown cells verbatim, code cells fenced with the
    notebook's language. A notebook is a document with code in it —
    general chunking treats it honestly."""
    import json

    nb = json.loads(file_bytes.decode("utf-8", errors="replace"))
    lang = (
        (nb.get("metadata") or {})
        .get("kernelspec", {})
        .get("language", "")
        or (nb.get("metadata") or {})
        .get("language_info", {})
        .get("name", "")
    )
    parts: list[str] = []
    for cell in nb.get("cells", []):
        source = cell.get("source") or []
        body = "".join(source) if isinstance(source, list) else str(source)
        body = body.rstrip()
        if not body.strip():
            continue
        kind = cell.get("cell_type")
        if kind == "markdown":
            parts.append(body)
        elif kind == "code":
            parts.append(f"```{lang}\n{body}\n```")
    return "\n\n".join(parts)


# --- Gate E (2026-07-20): tabular extraction -------------------------------
# One CANONICAL text form for every tabular source, so the chunker and
# the mechanical row extractor parse exactly one format: optional
# "=== SHEET: <name> ===" markers (xlsx only), first line = headers,
# tab-separated rows. Zero LLM anywhere in this lane (F4).

TABULAR_SHEET_MARKER = "=== SHEET: "


def extract_tabular_from_delimited(file_bytes: bytes, delimiter: str) -> str:
    """csv/tsv -> canonical TSV text (quotes and embedded delimiters
    resolved by the csv parser, tabs/newlines inside cells flattened)."""
    import csv
    import io
    text = file_bytes.decode("utf-8", errors="replace")
    out_lines = []
    for row in csv.reader(io.StringIO(text), delimiter=delimiter):
        out_lines.append("\t".join(
            (cell or "").replace("\t", " ").replace("\n", " ").strip()
            for cell in row
        ))
    return "\n".join(out_lines)


def extract_tabular_from_xlsx(file_bytes: bytes) -> str:
    """xlsx -> canonical text with sheet markers. read_only mode keeps
    memory flat on big workbooks; formulas arrive as computed values
    when the file carries them."""
    import io
    from openpyxl import load_workbook
    wb = load_workbook(
        io.BytesIO(file_bytes), read_only=True, data_only=True,
    )
    sections = []
    for ws in wb.worksheets:
        lines = [f"{TABULAR_SHEET_MARKER}{ws.title} ==="]
        for row in ws.iter_rows(values_only=True):
            cells = [
                str(c).replace("\t", " ").replace("\n", " ").strip()
                if c is not None else ""
                for c in row
            ]
            if any(cells):
                lines.append("\t".join(cells))
        if len(lines) > 1:
            sections.append("\n".join(lines))
    wb.close()
    return "\n\n".join(sections)


# --- Gate F (2026-07-20): conversational extraction ------------------------
# Canonical chat text: optional "=== WINDOW: YYYY-MM ===" markers
# (archives only, F-Q1=C), one unit (email / thread) per
# "--- <unit-ref> ---" section, speaker-attributed lines. Message-IDs
# and reply refs are PRESERVED in unit headers — slice 2's chain
# material.

CHAT_WINDOW_MARKER = "=== WINDOW: "
CHAT_UNIT_MARKER = "--- "


def _email_month(msg) -> str:
    from email.utils import parsedate_to_datetime
    try:
        dt = parsedate_to_datetime(msg.get("Date", ""))
        return f"{dt.year:04d}-{dt.month:02d}"
    except Exception:  # noqa: BLE001
        return "undated"


def _email_unit_text(msg) -> str:
    """One email -> one canonical unit."""
    def _body(m) -> str:
        if m.is_multipart():
            for part in m.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        return payload.decode(
                            part.get_content_charset() or "utf-8",
                            errors="replace",
                        )
            return ""
        payload = m.get_payload(decode=True)
        if payload:
            return payload.decode(
                m.get_content_charset() or "utf-8", errors="replace",
            )
        return str(m.get_payload() or "")

    refs = msg.get("In-Reply-To") or msg.get("References") or ""
    header = (
        f"{CHAT_UNIT_MARKER}message-id={msg.get('Message-ID', '').strip()}"
        f" in-reply-to={refs.strip().split()[-1] if refs.strip() else ''} ---"
    )
    lines = [
        header,
        f"From: {msg.get('From', '')} | Date: {msg.get('Date', '')}"
        f" | Subject: {msg.get('Subject', '')}",
    ]
    body = _body(msg).strip()
    sender = (msg.get("From") or "unknown").split("<")[0].strip() or "unknown"
    for ln in body.splitlines():
        if ln.strip():
            lines.append(f"{sender}: {ln.strip()}")
    return "\n".join(lines)


def extract_chat_from_eml(file_bytes: bytes) -> str:
    """Single email: whole-file, no window markers (F-Q1=C)."""
    import email
    from email import policy
    msg = email.message_from_bytes(file_bytes, policy=policy.default)
    return _email_unit_text(msg)


def extract_chat_from_mbox(file_bytes: bytes) -> str:
    """Archive: monthly windows (F-Q1=C — C4's worked example)."""
    import email
    import mailbox
    import os
    import tempfile
    from email import policy

    with tempfile.NamedTemporaryFile(
        suffix=".mbox", delete=False,
    ) as tmp:
        tmp.write(file_bytes)
        path = tmp.name
    try:
        box = mailbox.mbox(path)
        by_month: dict[str, list[str]] = {}
        for raw in box:
            msg = email.message_from_bytes(
                raw.as_bytes(), policy=policy.default,
            )
            by_month.setdefault(_email_month(msg), []).append(
                _email_unit_text(msg)
            )
        box.close()
    finally:
        os.unlink(path)

    sections = []
    for month in sorted(by_month):
        units = "\n\n".join(by_month[month])
        sections.append(f"{CHAT_WINDOW_MARKER}{month} ===\n{units}")
    return "\n\n".join(sections)


def looks_like_slack_export(text: str) -> bool:
    """Mechanical shape check: a JSON array of message objects with
    ts + (text|user). Generic JSON stays untouched — Gate G's
    schema-inference owns it."""
    import json
    head = text.lstrip()
    if not head.startswith("["):
        return False
    try:
        data = json.loads(text)
    except Exception:  # noqa: BLE001
        return False
    if not isinstance(data, list) or not data:
        return False
    sample = [d for d in data[:5] if isinstance(d, dict)]
    return bool(sample) and all(
        "ts" in d and ("text" in d or "user" in d) for d in sample
    )


def extract_chat_from_slack_json(text: str) -> str:
    """Slack channel export -> threads as units, monthly windows when
    the export spans months (F-Q1=C)."""
    import json
    from datetime import datetime, timezone

    data = json.loads(text)
    threads: dict[str, list[dict]] = {}
    order: list[str] = []
    for m in data:
        if not isinstance(m, dict):
            continue
        root = str(m.get("thread_ts") or m.get("ts") or "")
        if root not in threads:
            threads[root] = []
            order.append(root)
        threads[root].append(m)

    def _month(ts: str) -> str:
        try:
            dt = datetime.fromtimestamp(float(ts), tz=timezone.utc)
            return f"{dt.year:04d}-{dt.month:02d}"
        except Exception:  # noqa: BLE001
            return "undated"

    by_month: dict[str, list[str]] = {}
    for root in order:
        msgs = sorted(threads[root], key=lambda m: float(m.get("ts") or 0))
        lines = [f"{CHAT_UNIT_MARKER}thread={root} ---"]
        for m in msgs:
            who = m.get("user") or m.get("username") or "unknown"
            txt = (m.get("text") or "").replace("\n", " ").strip()
            if txt:
                lines.append(f"{who}: {txt}")
        if len(lines) > 1:
            by_month.setdefault(_month(root), []).append("\n".join(lines))

    if len(by_month) <= 1:
        # Single month / small export: whole-file (F-Q1=C).
        return "\n\n".join(
            u for units in by_month.values() for u in units
        )
    sections = []
    for month in sorted(by_month):
        units = "\n\n".join(by_month[month])
        sections.append(f"{CHAT_WINDOW_MARKER}{month} ===\n{units}")
    return "\n\n".join(sections)
